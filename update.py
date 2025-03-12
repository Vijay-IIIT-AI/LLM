from pptx import Presentation
import os

def is_empty_slide(slide):
    """Check if a slide is empty (no text, shapes, tables, or images)."""
    if not slide.shapes:
        return True
    for shape in slide.shapes:
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            if shape.text_frame and any(para.text.strip() for para in shape.text_frame.paragraphs):
                return False
        if hasattr(shape, "has_table") and shape.has_table:
            return False
        if hasattr(shape, "chart") and shape.chart:
            return False
        if hasattr(shape, "image") and shape.image:
            return False
    return True

def validate_ppt(file_path):
    """Validate a PowerPoint file based on the given conditions."""
    if not file_path.endswith(".pptx"):
        return "Invalid file type. Please provide a .pptx file."
    
    if not os.path.exists(file_path) or os.path.getsize(file_path) == 0:
        return "Invalid file: File is missing or empty."
    
    prs = Presentation(file_path)
    total_slides = len(prs.slides)
    empty_slides = [i+1 for i, slide in enumerate(prs.slides) if is_empty_slide(slide)]
    
    if total_slides > 20:
        return "Invalid file: More than 20 slides."
    
    if total_slides == 0 or len(empty_slides) == total_slides:
        return "Invalid file: All slides are empty."
    
    if len(empty_slides) > 0:
        return f"Warning: Empty slides found at positions {empty_slides}."
    
    return "Valid PowerPoint file."

# Example usage
file_path = r"C:\Users\Vijay\Desktop\PPT\New Microsoft PowerPoint Presentation.pptx"  # Change this to the actual file path
result = validate_ppt(file_path)
print(result)
